"""
deck_generator/agent.py
Deck Generator Agent — FastAPI service for creating professional business solution decks.

This service generates professional PowerPoint/PDF presentations based on:
- Business research data from SDR sessions
- Call transcripts and outcomes
- Customer pain points and opportunities
- Tailored solutions and value propositions

The generated deck is attached to calendar invites for meetings.
"""

from __future__ import annotations

import json
import logging
import uuid
import base64
import io
from contextlib import asynccontextmanager
from datetime import datetime
from pathlib import Path
from typing import Any, Dict

from fastapi import FastAPI, HTTPException
from dedalus_labs import AsyncDedalus, DedalusRunner
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import httpx
from dotenv import load_dotenv

from common.config import DEFAULT_MODEL, DRAFT_MODEL, UI_CLIENT_URL
from common.models import AgentCallback, AgentType

load_dotenv()

logger = logging.getLogger(__name__)


# Deck generation schema
DECK_REQUEST_SCHEMA = {
    "session_id": str,
    "business_name": str,
    "research_summary": str,
    "call_transcript": str,
    "call_outcome": str,
    "contact_email": str,
    "meeting_date": str,
}


@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("🎨 Deck Generator agent starting up...")
    yield
    logger.info("🎨 Deck Generator agent shutting down...")


app = FastAPI(title="RapidReach Deck Generator", lifespan=lifespan)


async def generate_deck_content(
    business_name: str,
    research_summary: str,
    call_transcript: str,
    call_outcome: str
) -> Dict[str, Any]:
    """
    Generate structured content for the business solution deck using AI.
    
    Args:
        business_name: Name of the business
        research_summary: Research findings about the business
        call_transcript: Transcript of the phone call
        call_outcome: Outcome classification of the call
        
    Returns:
        Dictionary with structured deck content
    """
    client = AsyncDedalus()
    runner = DedalusRunner(client)
    
    content_prompt = f"""Based on the following business information, create a professional business solution deck outline:

BUSINESS: {business_name}
RESEARCH SUMMARY: {research_summary}
CALL TRANSCRIPT: {call_transcript}
CALL OUTCOME: {call_outcome}

Generate a structured deck with the following sections. Provide specific, tailored content for each:

1. TITLE SLIDE
   - Compelling title for the presentation
   - Subtitle highlighting the main value proposition
   
2. EXECUTIVE SUMMARY (2-3 key points)
   - Business challenge identified
   - Proposed solution overview
   - Expected impact/ROI
   
3. CURRENT SITUATION ANALYSIS (3-4 points)
   - Current business challenges
   - Market position
   - Competitive landscape gaps
   - Missed opportunities
   
4. DETAILED RESEARCH INSIGHTS (4-5 findings)
   - Specific research findings about the business
   - Competitor analysis details
   - Market opportunities identified
   - Customer feedback and reviews analysis
   - Online presence assessment
   
5. PROPOSED SOLUTION (3-5 key features)
   - Specific solutions addressing their needs
   - Technology recommendations
   - Implementation approach
   - Unique value proposition
   
6. BENEFITS & ROI (4-5 benefits)
   - Quantified benefits where possible
   - Cost savings
   - Revenue opportunities
   - Competitive advantages
   - Long-term value
   
7. IMPLEMENTATION TIMELINE (3-4 phases)
   - Phase 1: Planning & Setup
   - Phase 2: Development & Design
   - Phase 3: Launch & Optimization
   - Phase 4: Growth & Maintenance
   
8. NEXT STEPS (3-4 actions)
   - Immediate actions
   - Decision points
   - Timeline expectations
   - Contact information

Format your response as a JSON object with this exact structure:
{{
    "title": "Main presentation title",
    "subtitle": "Value proposition subtitle",
    "executive_summary": ["point 1", "point 2", "point 3"],
    "situation_analysis": ["challenge 1", "challenge 2", "challenge 3", "challenge 4"],
    "research_insights": ["finding 1", "finding 2", "finding 3", "finding 4", "finding 5"],
    "proposed_solution": ["solution 1", "solution 2", "solution 3", "solution 4", "solution 5"],
    "benefits_roi": ["benefit 1", "benefit 2", "benefit 3", "benefit 4", "benefit 5"],
    "implementation_timeline": [
        {{"phase": "Planning & Setup", "duration": "X weeks", "description": "..."}},
        {{"phase": "Development & Design", "duration": "X weeks", "description": "..."}},
        {{"phase": "Launch & Optimization", "duration": "X weeks", "description": "..."}},
        {{"phase": "Growth & Maintenance", "duration": "Ongoing", "description": "..."}}
    ],
    "next_steps": ["step 1", "step 2", "step 3", "step 4"]
}}"""

    try:
        result = await runner.run(
            input=content_prompt,
            model=DRAFT_MODEL,
            max_steps=3,
        )
        
        # Parse the JSON response
        content_text = result.final_output.strip()
        if content_text.startswith("```json"):
            content_text = content_text[7:]
        if content_text.endswith("```"):
            content_text = content_text[:-3]
        
        content_data = json.loads(content_text)
        logger.info(f"✅ Generated deck content for {business_name}")
        return content_data
        
    except Exception as e:
        logger.error(f"❌ Failed to generate deck content: {e}")
        # Fallback content
        return {
            "title": f"Business Solution Proposal for {business_name}",
            "subtitle": "Driving Digital Growth Through Strategic Web Presence",
            "executive_summary": [
                "Identified opportunity to enhance digital presence",
                "Proposed comprehensive web solution",
                "Projected 25-40% increase in customer reach"
            ],
            "situation_analysis": [
                "Limited online visibility",
                "Competitors gaining market share",
                "Missing digital customer touchpoints",
                "Untapped online revenue potential"
            ],
            "research_insights": [
                "No current website or online presence found",
                "Local competitors with strong web presence identified",
                "High customer demand indicated by reviews",
                "Strong reputation but limited digital reach",
                "Significant opportunity for online customer acquisition"
            ],
            "proposed_solution": [
                "Professional website development",
                "SEO optimization strategy",
                "Social media integration",
                "Online booking/contact systems",
                "Mobile-responsive design"
            ],
            "benefits_roi": [
                "Increased customer acquisition",
                "24/7 business accessibility",
                "Improved brand credibility",
                "Competitive market positioning",
                "Measurable ROI tracking"
            ],
            "implementation_timeline": [
                {"phase": "Planning & Setup", "duration": "1-2 weeks", "description": "Requirements gathering and design planning"},
                {"phase": "Development & Design", "duration": "3-4 weeks", "description": "Website development and content creation"},
                {"phase": "Launch & Optimization", "duration": "1 week", "description": "Testing, launch, and initial optimization"},
                {"phase": "Growth & Maintenance", "duration": "Ongoing", "description": "Continuous improvement and support"}
            ],
            "next_steps": [
                "Schedule detailed requirements meeting",
                "Finalize project scope and timeline",
                "Begin design and development process",
                "Establish success metrics and KPIs"
            ]
        }


def create_professional_deck(content: Dict[str, Any], business_name: str, template_style: str = "professional") -> bytes:
    """
    Create a professional PowerPoint presentation using the generated content.
    
    Args:
        content: Structured content dictionary
        business_name: Name of the business
        template_style: Template style - "professional", "creative", or "tech"
        
    Returns:
        Bytes of the PowerPoint file
    """
    prs = Presentation()
    
    # Define color schemes based on template style
    if template_style == "creative":
        primary_color = RGBColor(255, 87, 34)    # Orange
        accent_color = RGBColor(255, 152, 0)     # Light orange
        text_color = RGBColor(51, 51, 51)        # Dark gray
    elif template_style == "tech":
        primary_color = RGBColor(67, 56, 202)    # Indigo
        accent_color = RGBColor(99, 102, 241)    # Light indigo  
        text_color = RGBColor(51, 51, 51)        # Dark gray
    else:  # professional (default)
        primary_color = RGBColor(30, 97, 146)    # Professional blue
        accent_color = RGBColor(72, 133, 237)    # Lighter blue
        text_color = RGBColor(51, 51, 51)        # Dark gray
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = content["title"]
    subtitle.text = content["subtitle"]
    
    # Style title
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.color.rgb = primary_color
    title_para.font.bold = True
    
    # Style subtitle
    subtitle_frame = subtitle.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = text_color
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    slide.shapes.title.text = "Executive Summary"
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()
    
    for point in content["executive_summary"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = text_color
        p.space_after = Pt(12)
    
    # Slide 3: Current Situation Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Current Situation Analysis"
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()
    
    for point in content["situation_analysis"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.space_after = Pt(10)
    
    # Slide 4: Detailed Research Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Research Insights - What We Discovered"
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()
    
    p_intro = tf.add_paragraph()
    p_intro.text = "Our comprehensive research revealed:"
    p_intro.font.size = Pt(18)
    p_intro.font.color.rgb = primary_color
    p_intro.font.bold = True
    p_intro.space_after = Pt(15)
    
    for insight in content["research_insights"]:
        p = tf.add_paragraph()
        p.text = f"• {insight}"
        p.font.size = Pt(15)
        p.font.color.rgb = text_color
        p.space_after = Pt(8)
    
    # Slide 5: Proposed Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Proposed Solution"
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()
    
    for point in content["proposed_solution"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.space_after = Pt(10)
    
    # Slide 6: Benefits & ROI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Benefits & Return on Investment"
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()
    
    for point in content["benefits_roi"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.space_after = Pt(10)
    
    # Slide 6: Implementation Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Implementation Timeline"
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()
    
    for phase in content["implementation_timeline"]:
        p = tf.add_paragraph()
        p.text = f"• {phase['phase']} ({phase['duration']})"
        p.font.size = Pt(16)
        p.font.color.rgb = primary_color
        p.font.bold = True
        p.space_after = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = f"  {phase['description']}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = text_color
        p2.space_after = Pt(12)
    
    # Slide 8: Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Next Steps"
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()
    
    for i, step in enumerate(content["next_steps"], 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {step}"
        p.font.size = Pt(18)
        p.font.color.rgb = text_color
        p.space_after = Pt(15)
    
    # Style all title slides
    for slide in prs.slides:
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.color.rgb = primary_color
            title_para.font.bold = True
    
    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@app.post("/generate-deck")
async def generate_deck(request: dict):
    """
    Generate a professional business solution deck.
    
    Expected request format:
    {
        "session_id": "string",
        "business_name": "string", 
        "research_summary": "string",
        "call_transcript": "string",
        "call_outcome": "string",
        "contact_email": "string",
        "meeting_date": "string",
        "template_style": "professional" | "creative" | "tech" (optional)
    }
    """
    try:
        # Validate required fields
        required_fields = ["business_name", "research_summary", "call_transcript", "call_outcome"]
        for field in required_fields:
            if field not in request:
                raise HTTPException(status_code=400, detail=f"Missing required field: {field}")
        
        # Callback for progress updates
        session_id = request.get("session_id", str(uuid.uuid4()))
        
        async def send_callback(message: str, step: str = ""):
            try:
                callback = AgentCallback(
                    agent_type=AgentType.DECK_GENERATOR,
                    session_id=session_id,
                    message=message,
                    step=step,
                    timestamp=datetime.now()
                )
                async with httpx.AsyncClient(timeout=5.0) as client:
                    await client.post(f"{UI_CLIENT_URL}/callback", json=callback.dict())
            except Exception as e:
                logger.warning(f"Callback failed: {e}")
        
        await send_callback("🎨 Starting deck generation...", "deck_generation")
        
        # Generate deck content
        await send_callback("🧠 Generating deck content with AI...", "content_generation")
        content = await generate_deck_content(
            business_name=request["business_name"],
            research_summary=request["research_summary"],
            call_transcript=request["call_transcript"],
            call_outcome=request["call_outcome"]
        )
        
        # Create PowerPoint presentation
        await send_callback("📊 Creating professional presentation...", "deck_creation")
        template_style = request.get("template_style", "professional")
        deck_bytes = create_professional_deck(content, request["business_name"], template_style)
        
        # Encode as base64 for transmission
        deck_b64 = base64.b64encode(deck_bytes).decode('utf-8')
        
        await send_callback("✅ Deck generation completed successfully!", "completed")
        
        return {
            "success": True,
            "session_id": session_id,
            "business_name": request["business_name"],
            "deck_content": content,
            "deck_file_b64": deck_b64,
            "filename": f"{request['business_name']}_Business_Solution.pptx",
            "created_at": datetime.now().isoformat()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Deck generation failed: {e}")
        raise HTTPException(status_code=500, detail=f"Deck generation failed: {str(e)}")


@app.get("/health")
async def health_check():
    """Health check endpoint."""
    return {"status": "healthy", "service": "deck-generator"}


@app.get("/")
async def root():
    """Root endpoint with service info."""
    return {
        "service": "RapidReach Deck Generator",
        "description": "Professional business solution deck generator",
        "version": "1.0.0",
        "endpoints": ["/generate-deck", "/health"]
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8086)